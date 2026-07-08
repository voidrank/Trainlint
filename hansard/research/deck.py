#!/usr/bin/env python3
"""Keynote-minimal slide deck for ANY hansard project — a permanent, substrate-driven
generator (the hand-built make_keynote.py, generalized).

  python3 deck.py [project]        # -> research/viz/<project>.keynote.pptx

It INVENTS nothing: every slide is folded from the same substrate the HTML report reads
(goal.<name>.txt + plan.<name>.jsonl + glossary), reusing viz._want_parts / plan.* so the
deck can never drift from the report. The DESIGN SYSTEM is the固化 part — one keynote
register (one statement per slide, huge type, one accent, consistent eyebrow+footer chrome,
data as a picture). Pure python-pptx: no browser, no node.

Archetypes, chosen by what the substrate carries:
  cover (goal headline + pillar chips) · the gates (pillars) · one claim+evidence slide per
  pillar · open-questions (open decisions, main thread starred) · closing (locked/open + bar).
"""
import re
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent
sys.path.insert(0, str(ROOT))
import plan          # noqa: E402
import viz           # noqa: E402  (reuse split_goal / _want_parts — single source w/ the report)
import paths         # noqa: E402  — per-project data lives outside the versioned plugin dir

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    sys.exit("deck.py needs python-pptx:  python3 -m pip install python-pptx")

# --- design tokens (the固化 register) -------------------------------------------------
BG   = RGBColor(0x0A, 0x0E, 0x1A); INK  = RGBColor(0xF2, 0xF6, 0xFC)
MUT  = RGBColor(0x7E, 0x8C, 0xA6); ACC  = RGBColor(0x4C, 0xC9, 0xF0)
ACC2 = RGBColor(0x2D, 0xD4, 0xBF); WARN = RGBColor(0xFB, 0xBF, 0x24)
OK   = RGBColor(0x34, 0xD3, 0x99); RULE = RGBColor(0x1E, 0x2A, 0x44)
FONT = "Arial"
MX = 1.0


def _trunc(s, n):
    s = " ".join(str(s or "").split())
    return s if len(s) <= n else s[:n - 1].rstrip() + "…"


def _takeaway(s, n=58):
    """First clause/sentence of a decision or choice — the punchy lead, not the whole para."""
    s = " ".join(str(s or "").split())
    m = re.search(r"[.;:](\s|$)", s)
    if m and m.start() < n * 1.6:
        s = s[:m.start()]
    return _trunc(s, n)


# --- LLM copy authoring (gen-time, via the Claude Code CLI) ----------------------------
# Turns the raw substrate into punchy, keynote-grade copy. Routed through `claude -p` headless
# (no API key — uses the CLI's own auth), so it works anywhere the CLI is installed. The model
# is told to COMPRESS, never invent. Returns an override map {cover_subtitle, closing,
# pillars:{id:{headline,accent,caption}}, open:{id:text}} or None (-> substrate-text fallback)
# when the CLI is missing or the call fails, so the deck always builds.
def _author_copy(name, goal, pillars, opens, mt):
    import os
    import json as _json
    import shutil
    import subprocess
    claude = shutil.which("claude")
    if not claude:
        sys.stderr.write("[deck] claude CLI not found; using substrate text\n")
        return None
    model = os.environ.get("DECK_MODEL", "sonnet")
    decs = [{"id": p.get("id"), "decision": p.get("decision", ""), "choice": p.get("choice", ""),
             "why": p.get("why", ""), "star": bool(mt and p.get("id") == mt.get("id"))}
            for p in pillars]
    op = [{"id": n.get("id"), "decision": n.get("decision", ""),
           "star": bool(mt and n.get("id") == mt.get("id"))} for n in opens]
    prompt = (
        "You write copy for a KEYNOTE-MINIMAL slide deck — one idea per slide, huge type, almost "
        "no words. Below is a research project's GOAL and key DECISIONS as JSON. COMPRESS them into "
        "punchy slide copy; never invent facts not present in the input. Output ONLY a JSON object "
        "— no prose, no code fence — matching exactly:\n"
        '{"cover_subtitle":"<=14 words","pillars":[{"id":"<id>","headline":"<=8 words, a claim not '
        'a question","accent":"the choice in <=12 words","caption":"why it matters <=16 words"}],'
        '"open":[{"id":"<id>","text":"<=9 words"}],"closing":"<=16 words"}\n\nPROJECT:\n'
        + _json.dumps({"project": name, "goal": goal, "pillars": decs, "open": op}, ensure_ascii=False))
    try:
        proc = subprocess.run([claude, "-p", prompt, "--output-format", "json", "--model", model],
                              capture_output=True, text=True, timeout=180)
        env = _json.loads(proc.stdout or "{}")
        text = env.get("result", "") if isinstance(env, dict) else ""
        m = re.search(r"\{.*\}", text, re.S)
        spec = _json.loads(m.group(0) if m else text)
    except Exception as e:
        sys.stderr.write(f"[deck] CLI copy authoring failed ({e}); falling back to substrate text\n")
        return None
    return {"cover_subtitle": spec.get("cover_subtitle", ""),
            "closing": spec.get("closing", ""),
            "pillars": {p.get("id"): p for p in spec.get("pillars", []) if p.get("id")},
            "open": {o.get("id"): o.get("text", "") for o in spec.get("open", []) if o.get("id")}}


class Deck:
    def __init__(self, name):
        self.name = name
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = Inches(13.333), Inches(7.5)
        self.blank = self.prs.slide_layouts[6]
        self.W, self.H = self.prs.slide_width, self.prs.slide_height
        self.builders = []   # deferred so footer N/total is correct
        self.total = 0

    # --- primitives ---
    def _slide(self, idx, section=""):
        s = self.prs.slides.add_slide(self.blank)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.W, self.H)
        bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background()
        bg.shadow.inherit = False
        if section:
            self._para(s, MX, 0.62, 11, 0.4, [(section.upper(), ACC2, True, 13)])
        ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MX), Inches(6.85), Inches(11.33), Pt(1))
        ln.fill.solid(); ln.fill.fore_color.rgb = RULE; ln.line.fill.background(); ln.shadow.inherit = False
        self._para(s, MX, 6.95, 11.33, 0.35, [(f"{self.name} · {idx}/{self.total}", MUT, False, 11)],
                   align=PP_ALIGN.RIGHT)
        return s

    def _para(self, s, x, y, w, h, items, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=10):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        for i, (t, c, b, sz) in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align; p.space_after = Pt(space)
            r = p.add_run(); r.text = t; r.font.size = Pt(sz); r.font.bold = b
            r.font.color.rgb = c; r.font.name = FONT
        return tb

    def _line(self, s, x, y, w, h, runs, size, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        p = tf.paragraphs[0]; p.alignment = align
        for t, c, b in runs:
            r = p.add_run(); r.text = t; r.font.size = Pt(size); r.font.bold = b
            r.font.color.rgb = c; r.font.name = FONT
        return tb

    def _bar(self, s, x, y, w, frac, track=RULE, fill=OK, hgt=0.18):
        t = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(hgt))
        t.fill.solid(); t.fill.fore_color.rgb = track; t.line.fill.background(); t.shadow.inherit = False
        if frac > 0:
            f = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                   Inches(max(0.06, w * frac)), Inches(hgt))
            f.fill.solid(); f.fill.fore_color.rgb = fill; f.line.fill.background(); f.shadow.inherit = False

    def add(self, fn):
        self.builders.append(fn)

    def save(self, path):
        self.total = len(self.builders)
        for i, fn in enumerate(self.builders, 1):
            fn(self._slide(i, fn.__doc__ or ""))
        self.prs.save(str(path))
        return path


def build(name):
    gp = paths.resolve(f"goal.{name}.txt")
    goal, bar = viz.split_goal(gp.read_text(encoding="utf-8") if gp.exists() else "")
    pl = plan.load(name)
    pillars = plan.pillars(pl)
    mt = plan.main_thread(pl)
    summ = plan.summary(pl); counts = summ["counts"]
    n_decided = counts.get("decided", 0) + counts.get("verified", 0)
    n_open = counts.get("open", 0)
    head, _bul, done = viz._want_parts(goal, bar, pl)
    opens = [n for n in pl if n.get("status") == "open"]
    spec = _author_copy(name, goal, pillars, opens, mt) or {}   # LLM copy; {} -> substrate text
    sp_pillars = spec.get("pillars", {})
    sp_open = spec.get("open", {})

    d = Deck(name)

    # cover
    def cover(s):
        d._line(s, MX, 2.4, 11.33, 1.4, [(name, INK, True)], 54, anchor=MSO_ANCHOR.MIDDLE)
        d._para(s, MX, 3.95, 10.6, 1.3, [(spec.get("cover_subtitle") or _trunc(head, 150), MUT, False, 22)])
        if pillars:
            chips = "     ·     ".join(p.get("id", "") for p in pillars)
            d._para(s, MX, 5.3, 11.33, 0.5, [(chips, ACC2, False, 14)])
    d.add(cover)

    # the gates (pillars)
    if pillars:
        def gates(s):
            """what we're building"""
            d._line(s, MX, 1.7, 11.33, 1.0,
                    [("%d pillars " % len(pillars), INK, True), ("hold this up.", ACC, True)],
                    44, anchor=MSO_ANCHOR.MIDDLE)
            items = [(f"◆  {sp_pillars.get(p.get('id'), {}).get('headline') or _takeaway(p.get('decision'), 72)}",
                      INK, False, 22) for p in pillars]
            d._para(s, MX, 3.15, 11.33, 3.0, items, space=10)
            if done:
                d._para(s, MX, 6.2, 11.33, 0.5, [("DONE = " + _trunc(done, 120), OK, False, 14)])
        d.add(gates)

    # one claim+evidence slide per pillar
    for p in pillars:
        star = " ★" if mt and p.get("id") == mt.get("id") else ""
        sect = f"pillar · {p.get('id','')}{star}"

        def pillar_slide(s, p=p, sect=sect):
            c = sp_pillars.get(p.get("id"), {})
            d._para(s, MX, 0.62, 11, 0.4, [(sect.upper(), ACC2, True, 13)])
            headline = c.get("headline") or _takeaway(p.get("decision"), 84)
            d._para(s, MX, 1.85, 11.33, 1.7, [(headline, INK, True, 40 if c.get("headline") else 36)], space=2)
            accent = c.get("accent") or (p.get("choice") and _trunc(p.get("choice"), 150))
            if accent:
                d._para(s, MX, 3.9, 11.33, 1.3, [("→ " + accent, ACC, False, 22)])
            caption = c.get("caption") or (p.get("why") and _trunc(p.get("why"), 170))
            if caption:
                d._para(s, MX, 5.7, 11.33, 0.9, [(caption, MUT, False, 18 if c.get("caption") else 17)])
        pillar_slide.__doc__ = ""   # section drawn inside (carries the ★)
        d.add(pillar_slide)

    # open questions
    if opens:
        def open_slide(s):
            """open questions"""
            d._line(s, MX, 1.7, 11.33, 1.0, [(str(len(opens)), ACC, True),
                    (" questions remain.", INK, True)], 50, anchor=MSO_ANCHOR.MIDDLE)
            items = []
            for n in opens:
                star = n.get("id") == (mt or {}).get("id")
                body = sp_open.get(n.get("id")) or _takeaway(n.get("decision"), 76)
                txt = ("★  " if star else "•  ") + body
                items.append((txt + ("   ← drive this next" if star else ""),
                              ACC if star else MUT, star, 22))
            d._para(s, MX, 3.2, 11.33, 3.0, items, space=8)
        d.add(open_slide)

    # closing
    def closing(s):
        d._line(s, MX, 2.4, 11.33, 1.0, [(f"{n_decided} locked", OK, True), ("   ·   ", MUT, False),
                (f"{n_open} open", WARN, True)], 46, anchor=MSO_ANCHOR.MIDDLE)
        total = max(1, summ["total"])
        d._bar(s, MX, 3.8, 11.33, n_decided / total)
        d._para(s, MX, 4.5, 11.33, 1.2, [(spec.get("closing") or _trunc(goal, 200), MUT, False, 18)])
    d.add(closing)

    outdir = ROOT / "viz"; outdir.mkdir(exist_ok=True)
    return d.save(outdir / f"{name}.keynote.pptx")


def main():
    import tree
    args = [a for a in sys.argv[1:] if not a.startswith("-")]
    name = tree._active(args[0] if args else None)
    path = build(name)
    print(f"keynote deck: {path}")


if __name__ == "__main__":
    main()
